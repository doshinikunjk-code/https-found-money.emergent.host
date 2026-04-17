from pathlib import Path
from dotenv import load_dotenv

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

from fastapi import FastAPI, APIRouter, HTTPException, Request, Response, UploadFile, File, Header, Query
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
from bson import ObjectId
from emergentintegrations.llm.chat import LlmChat, UserMessage
from emergentintegrations.payments.stripe.checkout import StripeCheckout, CheckoutSessionRequest
import requests as http_requests
import sib_api_v3_sdk
from sib_api_v3_sdk.rest import ApiException as BrevoApiException
import os
import logging
import bcrypt
import jwt
import uuid
import io
import json
import re
from datetime import datetime, timezone, timedelta
from pydantic import BaseModel
from typing import Optional, List
from fpdf import FPDF
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# MongoDB
mongo_url = os.environ.get('MONGO_URL', 'mongodb://localhost:27017')
client = AsyncIOMotorClient(mongo_url, tls=True, tlsAllowInvalidCertificates=True) if 'mongodb+srv' in mongo_url else AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

app = FastAPI()
api_router = APIRouter(prefix="/api")

JWT_ALGORITHM = "HS256"
EMERGENT_KEY = os.environ.get("EMERGENT_LLM_KEY")
STRIPE_KEY = os.environ.get("STRIPE_API_KEY")
BREVO_KEY = os.environ.get("BREVO_API_KEY")
SENDER_EMAIL = os.environ.get("NOTIFICATION_SENDER_EMAIL", "doshinikunjk@gmail.com")
SENDER_NAME = os.environ.get("NOTIFICATION_SENDER_NAME", "GrantGrabber")

# --- Brevo Email Notifications ---
def get_brevo_api():
    config = sib_api_v3_sdk.Configuration()
    config.api_key['api-key'] = BREVO_KEY
    return sib_api_v3_sdk.TransactionalEmailsApi(sib_api_v3_sdk.ApiClient(config))

async def send_notification_email(to_email: str, to_name: str, subject: str, html_content: str):
    if not BREVO_KEY:
        logger.warning("BREVO_API_KEY not set, skipping email")
        return
    try:
        api = get_brevo_api()
        email = sib_api_v3_sdk.SendSmtpEmail(
            sender={"email": SENDER_EMAIL, "name": SENDER_NAME},
            to=[{"email": to_email, "name": to_name}],
            subject=subject,
            html_content=html_content,
        )
        api.send_transac_email(email)
        logger.info(f"Email sent to {to_email}: {subject}")
    except BrevoApiException as e:
        logger.error(f"Brevo email error: {e}")
    except Exception as e:
        logger.error(f"Email send error: {e}")

def build_status_email(user_name, business_name, grant_name, old_status, new_status, grant_amount):
    status_labels = {
        "identified": "Identified", "draft_ready": "Draft Ready", "documents_required": "Documents Required",
        "under_review": "Under Review", "applied": "Applied", "submitted": "Submitted",
        "documents_submitted": "Documents Submitted", "approved": "Approved!", "rejected": "Rejected",
    }
    new_label = status_labels.get(new_status, new_status)
    color = "#10B981" if new_status in ["approved", "submitted", "applied"] else "#F59E0B" if new_status in ["draft_ready", "under_review"] else "#EF4444" if new_status == "rejected" else "#64748B"
    return f"""
    <div style="font-family:'Manrope',Arial,sans-serif;max-width:560px;margin:0 auto;padding:32px;background:#FAFAFA;">
      <div style="background:white;border-radius:12px;border:1px solid #E2E8F0;padding:32px;box-shadow:0 1px 3px rgba(0,0,0,0.05);">
        <div style="display:flex;align-items:center;gap:8px;margin-bottom:24px;">
          <div style="width:32px;height:32px;background:#10B981;border-radius:8px;display:inline-block;"></div>
          <span style="font-family:'Outfit',Arial,sans-serif;font-weight:700;font-size:18px;color:#0F172A;">GrantGrabber</span>
        </div>
        <h2 style="font-family:'Outfit',Arial,sans-serif;font-size:20px;color:#0F172A;margin:0 0 8px;">Grant Status Update</h2>
        <p style="color:#64748B;font-size:14px;margin:0 0 20px;">Hi {user_name},</p>
        <div style="background:#F8FAFC;border:1px solid #E2E8F0;border-radius:8px;padding:16px;margin-bottom:20px;">
          <p style="margin:0 0 4px;font-size:13px;color:#64748B;">Grant Program</p>
          <p style="margin:0 0 12px;font-size:16px;font-weight:600;color:#0F172A;">{grant_name}</p>
          <p style="margin:0 0 4px;font-size:13px;color:#64748B;">Business</p>
          <p style="margin:0 0 12px;font-size:14px;color:#0F172A;">{business_name}</p>
          <p style="margin:0 0 4px;font-size:13px;color:#64748B;">New Status</p>
          <span style="display:inline-block;background:{color};color:white;padding:4px 12px;border-radius:20px;font-size:13px;font-weight:600;">{new_label}</span>
        </div>
        <p style="color:#64748B;font-size:13px;margin:0;">Amount: <strong>${grant_amount:,.0f}</strong> | Log in to your dashboard for details.</p>
      </div>
    </div>"""

# Object Storage
STORAGE_URL = "https://integrations.emergentagent.com/objstore/api/v1/storage"
APP_NAME = "granthunter"
storage_key = None

def init_storage():
    global storage_key
    if storage_key:
        return storage_key
    resp = http_requests.post(f"{STORAGE_URL}/init", json={"emergent_key": EMERGENT_KEY}, timeout=30)
    resp.raise_for_status()
    storage_key = resp.json()["storage_key"]
    return storage_key

def put_object(path: str, data: bytes, content_type: str) -> dict:
    key = init_storage()
    resp = http_requests.put(
        f"{STORAGE_URL}/objects/{path}",
        headers={"X-Storage-Key": key, "Content-Type": content_type},
        data=data, timeout=120
    )
    resp.raise_for_status()
    return resp.json()

def get_object(path: str):
    key = init_storage()
    resp = http_requests.get(f"{STORAGE_URL}/objects/{path}", headers={"X-Storage-Key": key}, timeout=60)
    resp.raise_for_status()
    return resp.content, resp.headers.get("Content-Type", "application/octet-stream")

# Auth helpers
def get_jwt_secret():
    return os.environ["JWT_SECRET"]

def hash_password(password: str) -> str:
    return bcrypt.hashpw(password.encode("utf-8"), bcrypt.gensalt()).decode("utf-8")

def verify_password(plain: str, hashed: str) -> bool:
    return bcrypt.checkpw(plain.encode("utf-8"), hashed.encode("utf-8"))

def create_access_token(user_id: str, email: str) -> str:
    payload = {"sub": user_id, "email": email, "exp": datetime.now(timezone.utc) + timedelta(hours=24), "type": "access"}
    return jwt.encode(payload, get_jwt_secret(), algorithm=JWT_ALGORITHM)

async def get_current_user(request: Request) -> dict:
    token = request.cookies.get("access_token")
    if not token:
        auth_header = request.headers.get("Authorization", "")
        if auth_header.startswith("Bearer "):
            token = auth_header[7:]
    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    try:
        payload = jwt.decode(token, get_jwt_secret(), algorithms=[JWT_ALGORITHM])
        if payload.get("type") != "access":
            raise HTTPException(status_code=401, detail="Invalid token type")
        user = await db.users.find_one({"_id": ObjectId(payload["sub"])})
        if not user:
            raise HTTPException(status_code=401, detail="User not found")
        user["_id"] = str(user["_id"])
        user.pop("password_hash", None)
        return user
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expired")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")

# --- Pydantic Models ---
class RegisterInput(BaseModel):
    email: str
    password: str
    name: str
    business_name: Optional[str] = ""
    business_description: Optional[str] = ""

class LoginInput(BaseModel):
    email: str
    password: str

class GrantCreate(BaseModel):
    name: str
    description: str
    max_amount: float
    category: str
    eligibility: str
    provider: str
    grant_type: str = "non_refundable"

class ActionInput(BaseModel):
    action: str

class AIGenerateInput(BaseModel):
    service_type: str  # "grant_narrative", "business_plan", "pitch_deck"
    grant_id: Optional[str] = None
    context: Optional[str] = ""

class PaymentInput(BaseModel):
    grant_id: str
    origin_url: str

class DeadlineInput(BaseModel):
    grant_id: str
    deadline_date: str  # ISO date string
    reminder_days_before: int = 7
    notes: Optional[str] = ""

# --- Auth Endpoints ---
@api_router.post("/auth/register")
async def register(input: RegisterInput, response: Response):
    email = input.email.lower().strip()
    existing = await db.users.find_one({"email": email})
    if existing:
        raise HTTPException(status_code=400, detail="Email already registered")
    user_doc = {
        "email": email,
        "password_hash": hash_password(input.password),
        "name": input.name,
        "business_name": input.business_name or "",
        "business_description": input.business_description or "",
        "role": "client",
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    result = await db.users.insert_one(user_doc)
    user_id = str(result.inserted_id)
    token = create_access_token(user_id, email)
    await match_grants_for_user(user_id, input.business_name or "", input.business_description or "")
    return {"token": token, "user": {"id": user_id, "email": email, "name": input.name, "role": "client", "business_name": input.business_name, "onboarding_completed": False}}

@api_router.post("/auth/login")
async def login(input: LoginInput, response: Response):
    email = input.email.lower().strip()
    user = await db.users.find_one({"email": email})
    if not user or not verify_password(input.password, user["password_hash"]):
        raise HTTPException(status_code=401, detail="Invalid email or password")
    user_id = str(user["_id"])
    token = create_access_token(user_id, email)
    return {"token": token, "user": {"id": user_id, "email": email, "name": user.get("name", ""), "role": user.get("role", "client"), "business_name": user.get("business_name", ""), "onboarding_completed": user.get("onboarding_completed", False)}}

@api_router.post("/auth/logout")
async def logout(response: Response):
    response.delete_cookie("access_token", path="/")
    return {"message": "Logged out"}

@api_router.get("/auth/me")
async def get_me(request: Request):
    user = await get_current_user(request)
    return {"id": user["_id"], "email": user["email"], "name": user.get("name", ""), "role": user.get("role", "client"), "business_name": user.get("business_name", ""), "business_description": user.get("business_description", ""), "industry": user.get("industry", ""), "onboarding_completed": user.get("onboarding_completed", False), "has_logo": bool(user.get("logo_path"))}

# --- Onboarding ---
class OnboardingInput(BaseModel):
    business_name: Optional[str] = None
    business_description: Optional[str] = None
    industry: Optional[str] = None

@api_router.post("/onboarding/update-profile")
async def onboarding_update_profile(body: OnboardingInput, request: Request):
    user = await get_current_user(request)
    updates = {}
    if body.business_name is not None:
        updates["business_name"] = body.business_name
    if body.business_description is not None:
        updates["business_description"] = body.business_description
    if body.industry is not None:
        updates["industry"] = body.industry
    if updates:
        updates["updated_at"] = datetime.now(timezone.utc).isoformat()
        await db.users.update_one({"_id": ObjectId(user["_id"])}, {"$set": updates})
    updated = await db.users.find_one({"_id": ObjectId(user["_id"])}, {"_id": 0, "password_hash": 0})
    updated.pop("_id", None)
    return updated

@api_router.post("/onboarding/complete")
async def onboarding_complete(request: Request):
    user = await get_current_user(request)
    await db.users.update_one({"_id": ObjectId(user["_id"])}, {"$set": {"onboarding_completed": True, "updated_at": datetime.now(timezone.utc).isoformat()}})
    # Re-match grants with updated profile
    u = await db.users.find_one({"_id": ObjectId(user["_id"])})
    existing = await db.client_grants.count_documents({"user_id": user["_id"]})
    if existing == 0:
        await match_grants_for_user(user["_id"], u.get("business_name", ""), u.get("business_description", ""))
    return {"message": "Onboarding completed", "onboarding_completed": True}

# --- Dashboard ---
@api_router.get("/dashboard")
async def get_dashboard(request: Request):
    user = await get_current_user(request)
    user_id = user["_id"]
    client_grants = await db.client_grants.find({"user_id": user_id}, {"_id": 0}).to_list(100)
    enriched = []
    total_potential = 0
    total_secured = 0
    for cg in client_grants:
        grant = await db.grants.find_one({"id": cg["grant_id"]}, {"_id": 0})
        if grant:
            enriched.append({**cg, "grant": grant})
            total_potential += grant.get("max_amount", 0)
            if cg.get("status") == "approved":
                total_secured += grant.get("max_amount", 0)
    return {
        "user": {"id": user_id, "email": user["email"], "name": user.get("name", ""), "business_name": user.get("business_name", "")},
        "grants": enriched,
        "total_potential": total_potential,
        "total_secured": total_secured,
        "total_grants": len(enriched),
    }

# --- Client Grants ---
@api_router.get("/client/grants")
async def get_client_grants(request: Request):
    user = await get_current_user(request)
    client_grants = await db.client_grants.find({"user_id": user["_id"]}, {"_id": 0}).to_list(100)
    enriched = []
    for cg in client_grants:
        grant = await db.grants.find_one({"id": cg["grant_id"]}, {"_id": 0})
        if grant:
            enriched.append({**cg, "grant": grant})
    return enriched

@api_router.post("/client/grants/{grant_id}/action")
async def client_grant_action(grant_id: str, body: ActionInput, request: Request):
    user = await get_current_user(request)
    cg = await db.client_grants.find_one({"user_id": user["_id"], "grant_id": grant_id})
    if not cg:
        raise HTTPException(status_code=404, detail="Grant not found in your portfolio")
    old_status = cg.get("status", "identified")
    status_map = {"apply": "applied", "review": "under_review", "approve": "submitted", "upload_document": "documents_submitted"}
    new_status = status_map.get(body.action)
    if not new_status:
        raise HTTPException(status_code=400, detail="Invalid action")
    await db.client_grants.update_one(
        {"user_id": user["_id"], "grant_id": grant_id},
        {"$set": {"status": new_status, "updated_at": datetime.now(timezone.utc).isoformat()}}
    )
    updated = await db.client_grants.find_one({"user_id": user["_id"], "grant_id": grant_id}, {"_id": 0})
    grant = await db.grants.find_one({"id": grant_id}, {"_id": 0})
    # Send email notification
    if grant and old_status != new_status:
        html = build_status_email(user.get("name", ""), user.get("business_name", ""), grant["name"], old_status, new_status, grant.get("max_amount", 0))
        await send_notification_email(user["email"], user.get("name", ""), f"Grant Update: {grant['name']} - Status Changed", html)
    return {**updated, "grant": grant}

# --- AI Generation Services ---
@api_router.post("/ai/generate")
async def ai_generate(body: AIGenerateInput, request: Request):
    user = await get_current_user(request)
    business_name = user.get("business_name", "Your Business")
    business_desc = user.get("business_description", "")

    grant_info = ""
    if body.grant_id:
        grant = await db.grants.find_one({"id": body.grant_id}, {"_id": 0})
        if grant:
            grant_info = f"\nGrant: {grant['name']}\nProvider: {grant['provider']}\nMax Amount: ${grant['max_amount']:,.0f}\nEligibility: {grant['eligibility']}\nDescription: {grant['description']}\nType: {grant.get('grant_type', 'non_refundable')}"

    prompts = {
        "grant_narrative": f"""Write a grant application narrative for a Canadian government agency.
Business: {business_name}. {business_desc}
{grant_info}
{f'Context: {body.context}' if body.context else ''}
Include sections: Executive Summary, Business Overview, Project Description & Grant Alignment, Market Impact, Timeline (12-month), Budget Allocation, Expected Outcomes, Why Fund This Business. 600 words, professional tone.""",

        "business_plan": f"""Write a business plan for: {business_name}. {business_desc}
{f'Context: {body.context}' if body.context else ''}
Include: Executive Summary, Company Description, Products & Services, Market Analysis (TAM/SAM/SOM, competitors), Marketing Strategy, Operations, Team, 5-Year Financial Projections (revenue, expenses, margins), Funding Requirements, SWOT Analysis. 1000 words, investor-ready with numbers.""",

        "pitch_deck": f"""Create a 12-slide pitch deck for: {business_name}. {business_desc}
{f'Context: {body.context}' if body.context else ''}
For each slide provide headline, 4 bullets, and 2-sentence speaker notes.
Slides: 1-Title, 2-Problem, 3-Solution, 4-Market Opportunity (TAM/SAM/SOM), 5-Product, 6-Business Model & Revenue, 7-Traction, 8-Competitive Advantage, 9-Team, 10-Financials, 11-The Ask (funding), 12-Contact. Data-driven, compelling.""",
    }

    prompt = prompts.get(body.service_type)
    if not prompt:
        raise HTTPException(status_code=400, detail="Invalid service type. Use: grant_narrative, business_plan, or pitch_deck")

    try:
        chat = LlmChat(
            api_key=EMERGENT_KEY,
            session_id=f"ai-gen-{user['_id']}-{uuid.uuid4()}",
            system_message="You are a professional business consultant and grant writing expert. Provide well-structured, compelling content with clear section headers."
        ).with_model("anthropic", "claude-haiku-4-5-20251001")

        response = await chat.send_message(UserMessage(text=prompt))

        # Store generation record
        await db.ai_generations.insert_one({
            "id": str(uuid.uuid4()),
            "user_id": user["_id"],
            "service_type": body.service_type,
            "grant_id": body.grant_id,
            "content": response,
            "created_at": datetime.now(timezone.utc).isoformat(),
        })

        # If grant_narrative, update the client_grant draft
        if body.service_type == "grant_narrative" and body.grant_id:
            await db.client_grants.update_one(
                {"user_id": user["_id"], "grant_id": body.grant_id},
                {"$set": {"draft_content": response, "updated_at": datetime.now(timezone.utc).isoformat()}}
            )

        return {"content": response, "service_type": body.service_type}
    except Exception as e:
        logger.error(f"AI generation error: {e}")
        raise HTTPException(status_code=500, detail=f"AI generation failed: {str(e)}")

# --- PDF Generation ---
from fpdf.enums import XPos, YPos

def safe(text):
    return text.encode("latin-1", "replace").decode("latin-1")

@api_router.post("/ai/generate-pdf")
async def generate_pdf(request: Request):
    user = await get_current_user(request)
    body = await request.json()
    content = body.get("content", "")
    title = body.get("title", "Business Plan")
    business_name = user.get("business_name", "")
    date_str = datetime.now(timezone.utc).strftime("%B %d, %Y")

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=25)

    # === COVER PAGE ===
    pdf.add_page()
    # Dark header block
    pdf.set_fill_color(15, 23, 42)
    pdf.rect(0, 0, 210, 100, "F")
    # Emerald accent line
    pdf.set_fill_color(16, 185, 129)
    pdf.rect(0, 100, 210, 3, "F")

    # Client logo on cover
    logo_path = user.get("logo_path")
    if logo_path:
        try:
            logo_data, logo_ct = get_object(logo_path)
            ext = "png" if "png" in logo_ct else "jpg"
            tmp_logo = f"/tmp/logo_{user['_id']}.{ext}"
            with open(tmp_logo, "wb") as f:
                f.write(logo_data)
            pdf.image(tmp_logo, x=20, y=20, h=20)
        except Exception:
            pass

    # Title on dark background
    pdf.set_xy(20, 45)
    pdf.set_font("Helvetica", "B", 28)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(170, 14, safe(title), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    pdf.set_xy(20, 65)
    pdf.set_font("Helvetica", "", 14)
    pdf.set_text_color(16, 185, 129)
    pdf.cell(170, 8, safe(business_name), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    pdf.set_xy(20, 78)
    pdf.set_font("Helvetica", "", 10)
    pdf.set_text_color(148, 163, 184)
    pdf.cell(170, 6, safe(f"Prepared by GrantGrabber  |  {date_str}"), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    # Confidential note
    pdf.set_xy(20, 115)
    pdf.set_font("Helvetica", "I", 9)
    pdf.set_text_color(148, 163, 184)
    pdf.multi_cell(170, 5, safe("CONFIDENTIAL - This document contains proprietary business information prepared for internal use and grant application purposes."), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    # === CONTENT PAGES ===
    pdf.add_page()

    # Header helper for content pages
    def draw_page_header():
        pdf.set_fill_color(15, 23, 42)
        pdf.rect(0, 0, 210, 8, "F")
        pdf.set_fill_color(16, 185, 129)
        pdf.rect(0, 8, 210, 1.5, "F")
        pdf.set_font("Helvetica", "B", 6.5)
        pdf.set_text_color(255, 255, 255)
        pdf.set_xy(10, 1.5)
        pdf.cell(90, 5, safe("GRANTHUNTER AI"), new_x=XPos.RIGHT)
        pdf.cell(90, 5, safe(business_name.upper()), align="R", new_x=XPos.LMARGIN, new_y=YPos.NEXT)
        pdf.set_y(16)

    draw_page_header()

    section_num = 0
    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            pdf.ln(3)
            continue

        # Check page break
        if pdf.get_y() > 265:
            pdf.add_page()
            draw_page_header()

        # Main heading (# or ##)
        if stripped.startswith("# ") or stripped.startswith("## "):
            is_h1 = stripped.startswith("# ") and not stripped.startswith("## ")
            heading_text = stripped.lstrip("# ").strip()

            if is_h1:
                section_num += 1
                pdf.ln(6)
                # Section number + green bar
                pdf.set_fill_color(16, 185, 129)
                pdf.rect(20, pdf.get_y(), 40, 0.8, "F")
                pdf.ln(4)
                pdf.set_font("Helvetica", "B", 16)
                pdf.set_text_color(15, 23, 42)
                pdf.multi_cell(170, 9, safe(heading_text), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                pdf.ln(2)
            else:
                pdf.ln(4)
                pdf.set_font("Helvetica", "B", 12)
                pdf.set_text_color(15, 23, 42)
                pdf.multi_cell(170, 7, safe(heading_text), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
                pdf.ln(1)

            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(51, 65, 85)

        # ### subheading
        elif stripped.startswith("### "):
            pdf.ln(3)
            pdf.set_font("Helvetica", "B", 10.5)
            pdf.set_text_color(30, 41, 59)
            pdf.multi_cell(170, 6, safe(stripped.lstrip("# ").strip()), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(51, 65, 85)

        # Bold text **...**
        elif stripped.startswith("**") and stripped.endswith("**"):
            pdf.ln(2)
            pdf.set_font("Helvetica", "B", 10.5)
            pdf.set_text_color(30, 41, 59)
            pdf.multi_cell(170, 6, safe(stripped.strip("* ")), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(51, 65, 85)

        # Horizontal rule ---
        elif stripped == "---" or stripped == "***":
            pdf.ln(3)
            pdf.set_draw_color(226, 232, 240)
            pdf.set_line_width(0.3)
            pdf.line(20, pdf.get_y(), 190, pdf.get_y())
            pdf.ln(3)

        # Numbered list
        elif len(stripped) > 2 and stripped[0].isdigit() and stripped[1] in ".)" :
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(51, 65, 85)
            pdf.set_x(25)
            pdf.multi_cell(160, 5.5, safe(stripped), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        # Bullet points
        elif stripped.startswith("- ") or stripped.startswith("* "):
            bullet_text = stripped[2:].strip()
            # Check if bullet has bold prefix like **Key:** value
            if "**" in bullet_text:
                parts = bullet_text.split("**")
                pdf.set_x(25)
                pdf.set_font("Helvetica", "", 10)
                pdf.set_text_color(16, 185, 129)
                pdf.cell(4, 5.5, safe(chr(8226)), new_x=XPos.RIGHT)
                pdf.set_text_color(51, 65, 85)
                combined = bullet_text.replace("**", "")
                pdf.multi_cell(151, 5.5, safe(combined), new_x=XPos.LMARGIN, new_y=YPos.NEXT)
            else:
                pdf.set_x(25)
                pdf.set_font("Helvetica", "", 10)
                pdf.set_text_color(16, 185, 129)
                pdf.cell(4, 5.5, safe(chr(8226)), new_x=XPos.RIGHT)
                pdf.set_text_color(51, 65, 85)
                pdf.multi_cell(151, 5.5, safe(bullet_text), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

        # Regular paragraph
        else:
            pdf.set_font("Helvetica", "", 10)
            pdf.set_text_color(51, 65, 85)
            pdf.multi_cell(170, 5.5, safe(stripped), new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    # Footer on all pages
    for i in range(1, pdf.page + 1):
        pdf.page = i
        pdf.set_y(-18)
        pdf.set_draw_color(226, 232, 240)
        pdf.set_line_width(0.2)
        pdf.line(20, pdf.get_y(), 190, pdf.get_y())
        pdf.ln(3)
        pdf.set_font("Helvetica", "", 7)
        pdf.set_text_color(148, 163, 184)
        pdf.cell(85, 5, safe(f"GrantGrabber  |  {business_name}"), new_x=XPos.RIGHT)
        pdf.cell(85, 5, safe(f"Page {i}"), align="R", new_x=XPos.LMARGIN, new_y=YPos.NEXT)

    buf = io.BytesIO()
    pdf.output(buf)
    buf.seek(0)
    return Response(content=buf.read(), media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{title.replace(" ", "_")}.pdf"'})

# --- PPT Generation ---
@api_router.post("/ai/generate-ppt")
async def generate_ppt(request: Request):
    user = await get_current_user(request)
    body = await request.json()
    content = body.get("content", "")
    title = body.get("title", "Pitch Deck")
    business_name = user.get("business_name", "")

    prs = PptxPresentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    EMERALD = RGBColor(16, 185, 129)
    DARK = RGBColor(15, 23, 42)
    MID = RGBColor(51, 65, 85)
    LIGHT_GRAY = RGBColor(148, 163, 184)
    WHITE = RGBColor(255, 255, 255)
    LIGHT_BG = RGBColor(248, 250, 252)

    def add_slide_number(slide, num, total):
        from pptx.util import Emu
        txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7), Inches(1), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{num}/{total}"
        p.font.size = Pt(10)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.RIGHT

    def add_styled_slide(slide_title, bullets, notes_text=""):
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)
        # White background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # Top accent bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = EMERALD
        bar.line.fill.background()

        # Left emerald sidebar
        sidebar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
        sidebar.fill.solid()
        sidebar.fill.fore_color.rgb = EMERALD
        sidebar.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_title
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = DARK

        # Divider line under title
        div = slide.shapes.add_shape(1, Inches(0.8), Inches(1.5), Inches(2), Inches(0.04))
        div.fill.solid()
        div.fill.fore_color.rgb = EMERALD
        div.line.fill.background()

        # Bullets
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(11.5), Inches(4.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            clean = bullet.strip().lstrip("-*").strip()
            # Handle bold markers
            clean = clean.replace("**", "")
            p.text = clean
            p.font.size = Pt(20)
            p.font.color.rgb = MID
            p.space_after = Pt(12)
            p.level = 0

        # Footer bar
        footer_bar = slide.shapes.add_shape(1, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3))
        footer_bar.fill.solid()
        footer_bar.fill.fore_color.rgb = LIGHT_BG
        footer_bar.line.fill.background()

        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.2), Inches(5), Inches(0.3))
        fp = footer_box.text_frame.paragraphs[0]
        fp.text = f"{business_name}  |  GrantGrabber"
        fp.font.size = Pt(9)
        fp.font.color.rgb = LIGHT_GRAY

        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    # === TITLE SLIDE ===
    title_slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = title_slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

    # Emerald accent blocks
    accent1 = title_slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    accent1.fill.solid()
    accent1.fill.fore_color.rgb = EMERALD
    accent1.line.fill.background()

    accent2 = title_slide.shapes.add_shape(1, Inches(0), Inches(3.5), Inches(13.333), Inches(0.06))
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = EMERALD
    accent2.line.fill.background()

    # Logo
    logo_path = user.get("logo_path")
    if logo_path:
        try:
            logo_data, logo_ct = get_object(logo_path)
            ext = "png" if "png" in logo_ct else "jpg"
            tmp_logo = f"/tmp/ppt_logo_{user['_id']}.{ext}"
            with open(tmp_logo, "wb") as f:
                f.write(logo_data)
            title_slide.shapes.add_picture(tmp_logo, Inches(1), Inches(0.8), height=Inches(1))
        except Exception:
            pass

    # Title text
    title_box = title_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    sub_box = title_slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(0.8))
    sf = sub_box.text_frame
    p = sf.paragraphs[0]
    p.text = business_name
    p.font.size = Pt(24)
    p.font.color.rgb = EMERALD

    # Date
    date_box = title_slide.shapes.add_textbox(Inches(1), Inches(4.6), Inches(11), Inches(0.5))
    df = date_box.text_frame
    p = df.paragraphs[0]
    p.text = f"Prepared by GrantGrabber  |  {datetime.now(timezone.utc).strftime('%B %Y')}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GRAY

    # Confidential
    conf_box = title_slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11), Inches(0.4))
    cf = conf_box.text_frame
    p = cf.paragraphs[0]
    p.text = "CONFIDENTIAL"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = LIGHT_GRAY

    # === PARSE CONTENT INTO SLIDES ===
    current_title = ""
    current_bullets = []
    current_notes = ""
    for line in content.split("\n"):
        stripped = line.strip()
        if not stripped:
            continue
        is_heading = (stripped.startswith("Slide ") and ":" in stripped) or \
                     (stripped.startswith("# ") or stripped.startswith("## ")) or \
                     (stripped.startswith("**Slide"))
        if is_heading:
            if current_title and current_bullets:
                add_styled_slide(current_title, current_bullets, current_notes)
            clean = stripped.lstrip("#* ").strip()
            if ":" in clean:
                clean = clean.split(":", 1)[1].strip().strip("*")
            current_title = clean
            current_bullets = []
            current_notes = ""
        elif stripped.lower().startswith("speaker note") or stripped.lower().startswith("notes:"):
            current_notes = stripped.split(":", 1)[1].strip() if ":" in stripped else ""
        elif stripped.startswith("- ") or stripped.startswith("* ") or stripped.startswith("  -"):
            current_bullets.append(stripped)
        elif current_title and not stripped.startswith("For each"):
            if len(current_bullets) < 6:
                current_bullets.append(stripped)

    if current_title and current_bullets:
        add_styled_slide(current_title, current_bullets, current_notes)

    # Add slide numbers
    total_slides = len(prs.slides)
    for i, slide in enumerate(prs.slides):
        if i > 0:
            add_slide_number(slide, i, total_slides - 1)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Response(content=buf.read(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{title.replace(" ", "_")}.pptx"'})

# --- Real-Time AI Grant Matching ---
@api_router.post("/ai/match-grants")
async def ai_match_grants(request: Request):
    user = await get_current_user(request)
    business_name = user.get("business_name", "")
    business_desc = user.get("business_description", "")
    if not business_desc:
        raise HTTPException(status_code=400, detail="Please update your business description first")

    grants = await db.grants.find({}, {"_id": 0}).to_list(100)
    grants_text = "\n".join([f"- {g['name']} (Category: {g['category']}, Max: ${g['max_amount']:,.0f}, Type: {g.get('grant_type','non_refundable')}): Eligibility: {g['eligibility']}. Description: {g['description']}" for g in grants])

    prompt = f"""Analyze this business against each grant program and provide match scores.

Business: {business_name}
Description: {business_desc}

Grant Programs:
{grants_text}

For EACH grant, provide a JSON object with:
- "name": grant name (exact match)
- "score": match score 0-100 (be realistic - consider industry relevance, eligibility fit, business stage)
- "status": one of "draft_ready" (score>=80), "identified" (score>=40), or "low_match" (score<40)
- "reasoning": one sentence explaining the score

Return ONLY a valid JSON array. No markdown, no explanation outside the array.
Example: [{{"name":"CanExport SME","score":85,"status":"draft_ready","reasoning":"Strong export potential for health-tech products."}}]"""

    try:
        chat = LlmChat(
            api_key=EMERGENT_KEY,
            session_id=f"match-{user['_id']}-{uuid.uuid4()}",
            system_message="You are a grant eligibility analyst. Return ONLY valid JSON arrays. Be realistic with scores."
        ).with_model("anthropic", "claude-haiku-4-5-20251001")

        response = await chat.send_message(UserMessage(text=prompt))

        # Parse JSON from response
        json_match = re.search(r'\[.*\]', response, re.DOTALL)
        if not json_match:
            raise ValueError("No JSON array found in AI response")
        matches = json.loads(json_match.group())

        grant_map = {g["name"]: g for g in grants}
        updated_grants = []
        for m in matches:
            grant = grant_map.get(m.get("name"))
            if not grant:
                continue
            score = max(0, min(100, int(m.get("score", 50))))
            status = m.get("status", "identified")
            if status not in ["draft_ready", "identified", "low_match"]:
                status = "draft_ready" if score >= 80 else "identified" if score >= 40 else "identified"

            existing_cg = await db.client_grants.find_one({"user_id": user["_id"], "grant_id": grant["id"]})
            if existing_cg:
                # Only update score/reasoning, preserve user-driven status
                update_fields = {"match_score": score, "ai_reasoning": m.get("reasoning", ""), "updated_at": datetime.now(timezone.utc).isoformat()}
                if existing_cg.get("status") in ["identified", "low_match"]:
                    update_fields["status"] = status
                await db.client_grants.update_one({"user_id": user["_id"], "grant_id": grant["id"]}, {"$set": update_fields})
            else:
                await db.client_grants.insert_one({
                    "id": str(uuid.uuid4()), "user_id": user["_id"], "grant_id": grant["id"],
                    "match_score": score, "status": status, "draft_content": "", "documents": [],
                    "notes": "", "ai_reasoning": m.get("reasoning", ""), "payment_status": "unpaid",
                    "created_at": datetime.now(timezone.utc).isoformat(), "updated_at": datetime.now(timezone.utc).isoformat(),
                })
            updated_grants.append({"name": grant["name"], "score": score, "status": status, "reasoning": m.get("reasoning", "")})

        # Send email about re-matching
        html_rows = "".join([f"<tr><td style='padding:6px 10px;border-bottom:1px solid #E2E8F0;font-size:13px;'>{g['name']}</td><td style='padding:6px 10px;border-bottom:1px solid #E2E8F0;font-size:13px;font-weight:600;color:#10B981;'>{g['score']}%</td><td style='padding:6px 10px;border-bottom:1px solid #E2E8F0;font-size:12px;color:#64748B;'>{g['reasoning']}</td></tr>" for g in updated_grants])
        email_html = f"""<div style="font-family:Arial,sans-serif;max-width:560px;margin:0 auto;padding:32px;background:#FAFAFA;">
          <div style="background:white;border-radius:12px;border:1px solid #E2E8F0;padding:24px;">
            <h2 style="font-size:18px;color:#0F172A;margin:0 0 8px;">AI Grant Matching Complete</h2>
            <p style="color:#64748B;font-size:13px;margin:0 0 16px;">Hi {user.get('name','')}, your grants have been re-analyzed by AI.</p>
            <table style="width:100%;border-collapse:collapse;"><thead><tr style="background:#F8FAFC;"><th style="padding:8px 10px;text-align:left;font-size:12px;color:#64748B;">Grant</th><th style="padding:8px 10px;text-align:left;font-size:12px;color:#64748B;">Score</th><th style="padding:8px 10px;text-align:left;font-size:12px;color:#64748B;">Reasoning</th></tr></thead><tbody>{html_rows}</tbody></table>
          </div></div>"""
        await send_notification_email(user["email"], user.get("name", ""), "GrantGrabber: Your Grant Matches Updated", email_html)

        return {"matches": updated_grants, "total": len(updated_grants)}
    except json.JSONDecodeError as e:
        logger.error(f"AI match JSON parse error: {e}")
        raise HTTPException(status_code=500, detail="AI returned invalid format. Please try again.")
    except Exception as e:
        logger.error(f"AI matching error: {e}")
        raise HTTPException(status_code=500, detail=f"AI matching failed: {str(e)}")

# --- File Upload ---
@api_router.post("/upload")
async def upload_file(request: Request, file: UploadFile = File(...)):
    user = await get_current_user(request)
    ext = file.filename.split(".")[-1] if "." in file.filename else "bin"
    path = f"{APP_NAME}/uploads/{user['_id']}/{uuid.uuid4()}.{ext}"
    data = await file.read()
    result = put_object(path, data, file.content_type or "application/octet-stream")
    file_doc = {
        "id": str(uuid.uuid4()),
        "user_id": user["_id"],
        "storage_path": result["path"],
        "original_filename": file.filename,
        "content_type": file.content_type,
        "size": result.get("size", len(data)),
        "is_deleted": False,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.files.insert_one(file_doc)
    file_doc.pop("_id", None)
    return file_doc

@api_router.post("/client/grants/{grant_id}/upload-doc")
async def upload_grant_document(grant_id: str, request: Request, file: UploadFile = File(...)):
    user = await get_current_user(request)
    cg = await db.client_grants.find_one({"user_id": user["_id"], "grant_id": grant_id})
    if not cg:
        raise HTTPException(status_code=404, detail="Grant not found in your portfolio")

    ext = file.filename.split(".")[-1] if "." in file.filename else "bin"
    path = f"{APP_NAME}/grants/{user['_id']}/{grant_id}/{uuid.uuid4()}.{ext}"
    data = await file.read()
    result = put_object(path, data, file.content_type or "application/octet-stream")

    file_ref = {
        "id": str(uuid.uuid4()),
        "storage_path": result["path"],
        "original_filename": file.filename,
        "content_type": file.content_type,
        "size": result.get("size", len(data)),
        "uploaded_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.client_grants.update_one(
        {"user_id": user["_id"], "grant_id": grant_id},
        {"$push": {"documents": file_ref}, "$set": {"status": "documents_submitted", "updated_at": datetime.now(timezone.utc).isoformat()}}
    )
    return {"message": "Document uploaded", "file": file_ref}

@api_router.get("/files/{path:path}")
async def download_file(path: str, request: Request, auth: str = Query(None)):
    token = None
    if auth:
        token = auth
    else:
        auth_header = request.headers.get("Authorization", "")
        if auth_header.startswith("Bearer "):
            token = auth_header[7:]
    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    record = await db.files.find_one({"storage_path": path, "is_deleted": False})
    if not record:
        raise HTTPException(status_code=404, detail="File not found")
    data, content_type = get_object(path)
    return Response(content=data, media_type=record.get("content_type", content_type))

# --- Stripe Payments ---
COMMISSION_RATE = 0.10

@api_router.post("/payments/create-checkout")
async def create_checkout(body: PaymentInput, request: Request):
    user = await get_current_user(request)
    cg = await db.client_grants.find_one({"user_id": user["_id"], "grant_id": body.grant_id})
    if not cg:
        raise HTTPException(status_code=404, detail="Grant not found")
    grant = await db.grants.find_one({"id": body.grant_id}, {"_id": 0})
    if not grant:
        raise HTTPException(status_code=404, detail="Grant details not found")

    amount = round(grant["max_amount"] * COMMISSION_RATE, 2)
    success_url = f"{body.origin_url}/dashboard?payment=success&session_id={{CHECKOUT_SESSION_ID}}"
    cancel_url = f"{body.origin_url}/dashboard?payment=cancelled"

    webhook_url = f"{str(request.base_url).rstrip('/')}/api/webhook/stripe"
    stripe_checkout = StripeCheckout(api_key=STRIPE_KEY, webhook_url=webhook_url)

    checkout_req = CheckoutSessionRequest(
        amount=float(amount),
        currency="cad",
        success_url=success_url,
        cancel_url=cancel_url,
        metadata={"user_id": user["_id"], "grant_id": body.grant_id, "grant_name": grant["name"], "commission_rate": str(COMMISSION_RATE)}
    )
    session = await stripe_checkout.create_checkout_session(checkout_req)

    await db.payment_transactions.insert_one({
        "id": str(uuid.uuid4()),
        "session_id": session.session_id,
        "user_id": user["_id"],
        "grant_id": body.grant_id,
        "amount": amount,
        "currency": "cad",
        "grant_name": grant["name"],
        "payment_status": "initiated",
        "created_at": datetime.now(timezone.utc).isoformat(),
    })
    return {"url": session.url, "session_id": session.session_id}

@api_router.get("/payments/status/{session_id}")
async def get_payment_status(session_id: str, request: Request):
    user = await get_current_user(request)
    tx = await db.payment_transactions.find_one({"session_id": session_id, "user_id": user["_id"]}, {"_id": 0})
    if not tx:
        raise HTTPException(status_code=404, detail="Transaction not found")

    if tx.get("payment_status") not in ["paid", "expired"]:
        try:
            webhook_url = f"{str(request.base_url).rstrip('/')}/api/webhook/stripe"
            stripe_checkout = StripeCheckout(api_key=STRIPE_KEY, webhook_url=webhook_url)
            status = await stripe_checkout.get_checkout_status(session_id)
            new_status = status.payment_status if status.payment_status else tx.get("payment_status", "initiated")
            await db.payment_transactions.update_one(
                {"session_id": session_id},
                {"$set": {"payment_status": new_status, "status": status.status, "updated_at": datetime.now(timezone.utc).isoformat()}}
            )
            if new_status == "paid":
                existing_paid = await db.payment_transactions.find_one({"session_id": session_id, "payment_status": "paid", "credited": True})
                if not existing_paid:
                    await db.payment_transactions.update_one({"session_id": session_id}, {"$set": {"credited": True}})
                    await db.client_grants.update_one(
                        {"user_id": user["_id"], "grant_id": tx["grant_id"]},
                        {"$set": {"payment_status": "paid", "updated_at": datetime.now(timezone.utc).isoformat()}}
                    )
            tx = await db.payment_transactions.find_one({"session_id": session_id}, {"_id": 0})
        except Exception as e:
            logger.error(f"Payment status check error: {e}")

    return tx

@api_router.post("/webhook/stripe")
async def stripe_webhook(request: Request):
    body = await request.body()
    sig = request.headers.get("Stripe-Signature")
    try:
        webhook_url = f"{str(request.base_url).rstrip('/')}/api/webhook/stripe"
        stripe_checkout = StripeCheckout(api_key=STRIPE_KEY, webhook_url=webhook_url)
        event = await stripe_checkout.handle_webhook(body, sig)
        if event and event.payment_status == "paid":
            await db.payment_transactions.update_one(
                {"session_id": event.session_id},
                {"$set": {"payment_status": "paid", "updated_at": datetime.now(timezone.utc).isoformat()}}
            )
    except Exception as e:
        logger.error(f"Webhook error: {e}")
    return {"status": "ok"}

# --- Client Logo/Branding ---
@api_router.post("/client/upload-logo")
async def upload_logo(request: Request, file: UploadFile = File(...)):
    user = await get_current_user(request)
    allowed_types = ["image/png", "image/jpeg", "image/jpg", "image/webp", "image/svg+xml"]
    if file.content_type not in allowed_types:
        raise HTTPException(status_code=400, detail="Only PNG, JPEG, WEBP, or SVG images allowed")
    ext = file.filename.split(".")[-1] if "." in file.filename else "png"
    path = f"{APP_NAME}/logos/{user['_id']}/{uuid.uuid4()}.{ext}"
    data = await file.read()
    if len(data) > 2 * 1024 * 1024:
        raise HTTPException(status_code=400, detail="Logo must be under 2MB")
    result = put_object(path, data, file.content_type or "image/png")
    await db.users.update_one({"_id": ObjectId(user["_id"])}, {"$set": {
        "logo_path": result["path"],
        "logo_content_type": file.content_type,
        "logo_original_name": file.filename,
        "updated_at": datetime.now(timezone.utc).isoformat(),
    }})
    return {"message": "Logo uploaded", "logo_path": result["path"]}

@api_router.get("/client/logo")
async def get_logo(request: Request):
    user = await get_current_user(request)
    logo_path = user.get("logo_path")
    if not logo_path:
        raise HTTPException(status_code=404, detail="No logo uploaded")
    data, ct = get_object(logo_path)
    return Response(content=data, media_type=user.get("logo_content_type", ct))

@api_router.get("/client/branding")
async def get_branding(request: Request):
    user = await get_current_user(request)
    return {
        "has_logo": bool(user.get("logo_path")),
        "business_name": user.get("business_name", ""),
        "logo_url": f"/api/client/logo" if user.get("logo_path") else None,
    }

# --- Grant Deadlines & Calendar ---
@api_router.post("/deadlines")
async def create_deadline(body: DeadlineInput, request: Request):
    user = await get_current_user(request)
    grant = await db.grants.find_one({"id": body.grant_id}, {"_id": 0})
    if not grant:
        raise HTTPException(status_code=404, detail="Grant not found")
    deadline_doc = {
        "id": str(uuid.uuid4()),
        "user_id": user["_id"],
        "grant_id": body.grant_id,
        "grant_name": grant["name"],
        "deadline_date": body.deadline_date,
        "reminder_days_before": body.reminder_days_before,
        "notes": body.notes or "",
        "reminder_sent": False,
        "created_at": datetime.now(timezone.utc).isoformat(),
    }
    await db.deadlines.insert_one(deadline_doc)
    deadline_doc.pop("_id", None)
    return deadline_doc

@api_router.get("/deadlines")
async def get_deadlines(request: Request):
    user = await get_current_user(request)
    deadlines = await db.deadlines.find({"user_id": user["_id"]}, {"_id": 0}).to_list(100)
    return deadlines

@api_router.delete("/deadlines/{deadline_id}")
async def delete_deadline(deadline_id: str, request: Request):
    user = await get_current_user(request)
    result = await db.deadlines.delete_one({"id": deadline_id, "user_id": user["_id"]})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Deadline not found")
    return {"message": "Deadline deleted"}

@api_router.post("/deadlines/check-reminders")
async def check_reminders(request: Request):
    """Check all deadlines and send email reminders for upcoming ones."""
    user = await get_current_user(request)
    now = datetime.now(timezone.utc)
    deadlines = await db.deadlines.find({"user_id": user["_id"], "reminder_sent": False}, {"_id": 0}).to_list(100)
    sent = 0
    for dl in deadlines:
        try:
            deadline_dt = datetime.fromisoformat(dl["deadline_date"].replace("Z", "+00:00"))
            if deadline_dt.tzinfo is None:
                deadline_dt = deadline_dt.replace(tzinfo=timezone.utc)
            days_until = (deadline_dt - now).days
            if 0 <= days_until <= dl.get("reminder_days_before", 7):
                html = f"""<div style="font-family:Arial,sans-serif;max-width:560px;margin:0 auto;padding:32px;background:#FAFAFA;">
                  <div style="background:white;border-radius:12px;border:1px solid #E2E8F0;padding:32px;">
                    <div style="margin-bottom:20px;"><span style="font-weight:700;font-size:18px;color:#0F172A;">GrantGrabber</span></div>
                    <h2 style="font-size:20px;color:#0F172A;margin:0 0 8px;">Deadline Reminder</h2>
                    <p style="color:#64748B;font-size:14px;">Hi {user.get('name','')},</p>
                    <div style="background:#FEF3C7;border:1px solid #FCD34D;border-radius:8px;padding:16px;margin:16px 0;">
                      <p style="margin:0 0 8px;font-size:16px;font-weight:600;color:#92400E;">{dl['grant_name']}</p>
                      <p style="margin:0;font-size:14px;color:#92400E;">Deadline: <strong>{deadline_dt.strftime('%B %d, %Y')}</strong> ({days_until} days away)</p>
                    </div>
                    {f'<p style="color:#64748B;font-size:13px;">Notes: {dl["notes"]}</p>' if dl.get('notes') else ''}
                    <p style="color:#64748B;font-size:13px;">Log in to your dashboard to take action before the deadline.</p>
                  </div></div>"""
                await send_notification_email(user["email"], user.get("name", ""), f"Deadline Reminder: {dl['grant_name']} - {days_until} days left", html)
                await db.deadlines.update_one({"id": dl["id"]}, {"$set": {"reminder_sent": True, "reminder_sent_at": now.isoformat()}})
                sent += 1
        except Exception as e:
            logger.error(f"Reminder check error for {dl.get('id')}: {e}")
    return {"reminders_sent": sent, "total_checked": len(deadlines)}

# --- Grants (public) ---
@api_router.get("/grants")
async def get_grants():
    return await db.grants.find({}, {"_id": 0}).to_list(100)

# --- Admin ---
@api_router.get("/admin/clients")
async def admin_get_clients(request: Request):
    user = await get_current_user(request)
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    clients = []
    async for u in db.users.find({"role": "client"}):
        uid = str(u["_id"])
        grant_count = await db.client_grants.count_documents({"user_id": uid})
        clients.append({"id": uid, "email": u["email"], "name": u.get("name", ""), "business_name": u.get("business_name", ""), "grant_count": grant_count, "created_at": u.get("created_at", "")})
    return clients

@api_router.get("/admin/stats")
async def admin_get_stats(request: Request):
    user = await get_current_user(request)
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    return {
        "total_clients": await db.users.count_documents({"role": "client"}),
        "total_grants": await db.grants.count_documents({}),
        "total_applications": await db.client_grants.count_documents({"status": {"$in": ["applied", "submitted", "approved"]}}),
        "total_approved": await db.client_grants.count_documents({"status": "approved"}),
    }

@api_router.post("/admin/grants")
async def admin_create_grant(grant: GrantCreate, request: Request):
    user = await get_current_user(request)
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    grant_doc = {"id": str(uuid.uuid4()), **grant.model_dump(), "created_at": datetime.now(timezone.utc).isoformat()}
    await db.grants.insert_one(grant_doc)
    grant_doc.pop("_id", None)
    return grant_doc

@api_router.put("/admin/client-grants/{cg_id}")
async def admin_update_cg(cg_id: str, request: Request):
    user = await get_current_user(request)
    if user.get("role") != "admin":
        raise HTTPException(status_code=403, detail="Admin access required")
    body = await request.json()
    allowed = {k: v for k, v in body.items() if k in ["status", "match_score", "draft_content", "notes"]}
    if allowed:
        allowed["updated_at"] = datetime.now(timezone.utc).isoformat()
        await db.client_grants.update_one({"id": cg_id}, {"$set": allowed})
    updated = await db.client_grants.find_one({"id": cg_id}, {"_id": 0})
    return updated

# --- Grant Matching ---
async def match_grants_for_user(user_id: str, business_name: str, business_desc: str):
    grants = await db.grants.find({}, {"_id": 0}).to_list(100)
    configs = {
        "Export": {"score": 96, "status": "draft_ready", "draft": f"Export Market Development Application\n\nBusiness: {business_name or 'Your Business'}\n\nThis application outlines the export readiness and international market development strategy. The business has demonstrated strong potential for international expansion.\n\nKey highlights:\n- Established domestic customer base\n- Scalable production capacity\n- Identified target international markets\n- Clear export development timeline"},
        "Workforce": {"score": 78, "status": "documents_required", "draft": ""},
        "Clean Tech": {"score": 65, "status": "identified", "draft": ""},
        "Digital": {"score": 88, "status": "draft_ready", "draft": f"Digital Adoption Strategy\n\nBusiness: {business_name or 'Your Business'}\n\nComprehensive digital transformation plan including e-commerce, digital marketing automation, and CRM implementation."},
        "Innovation": {"score": 45, "status": "identified", "draft": ""},
        "Health Tech": {"score": 92, "status": "draft_ready", "draft": f"Health Technology Innovation Application\n\nBusiness: {business_name or 'Your Business'}\n\nThis application details how the business leverages technology to address healthcare challenges, with a focus on patient outcomes and accessibility."},
    }
    for grant in grants:
        cat = grant.get("category", "")
        cfg = configs.get(cat, {"score": 50, "status": "identified", "draft": ""})
        cg = {
            "id": str(uuid.uuid4()), "user_id": user_id, "grant_id": grant["id"],
            "match_score": cfg["score"], "status": cfg["status"], "draft_content": cfg["draft"],
            "documents": [], "notes": "", "payment_status": "unpaid",
            "created_at": datetime.now(timezone.utc).isoformat(), "updated_at": datetime.now(timezone.utc).isoformat(),
        }
        await db.client_grants.insert_one(cg)

# --- Seed Data ---
SEED_GRANTS = [
    {"name": "CanExport SME", "description": "Export development funding for Canadian SMEs expanding into international markets. Covers market research, trade shows, and export strategy.", "max_amount": 50000, "category": "Export", "eligibility": "Canadian SME with export potential, <500 employees, revenue under $100M", "provider": "Government of Canada", "grant_type": "non_refundable",
     "agency_details": {"portal_url": "https://www.tradecommissioner.gc.ca/en/our-solutions/funding-financing-international-business/canexport-smes/how-to-apply.html", "portal_name": "Trade Commissioner Service Portal", "contact_phone": "1-888-306-9991", "contact_email": "canexport-sme@international.gc.ca", "deadline": "May 29, 2026 at 12:00 PM ET", "processing_time": "60-90 business days", "how_to_apply": "1. Register at Trade Commissioner Service portal\n2. Complete My Profile and Account Profile\n3. Upload GST 34 return for eligibility verification\n4. Submit under My Applications > New Application > SME\n5. Apply at least 60 business days before activity start date"}},
    {"name": "Canada-Ontario Job Grant", "description": "Training grant to help employers invest in workforce development. Covers up to two-thirds of training costs for existing or new employees.", "max_amount": 10000, "category": "Workforce", "eligibility": "Ontario-based employer with employee training plan, max 25 participants per application", "provider": "Ontario Government", "grant_type": "non_refundable",
     "agency_details": {"portal_url": "https://eopg.labour.gov.on.ca/en/programs/canada-ontario-job-grant/application-process/", "portal_name": "Employment Ontario Partners' Gateway (EOPG)", "contact_phone": "1-800-387-5656", "contact_email": "EmploymentHotlineInquiries@ontario.ca", "deadline": "Rolling intake - apply anytime", "processing_time": "4-8 weeks", "how_to_apply": "1. Register at Employment Ontario Partners' Gateway (EOPG)\n2. Complete the COJG Employer Application form\n3. Attach Training Agreement (Form 2953E)\n4. Include Participant Registration (Form 2945E)\n5. Submit online through the portal"}},
    {"name": "Clean Technology Grant (EIP)", "description": "Energy Innovation Program funding for businesses adopting clean technology solutions. Supports energy efficiency, clean fuels, renewables, and sustainable practices.", "max_amount": 75000, "category": "Clean Tech", "eligibility": "Canadian-incorporated entity implementing clean tech solutions, co-funding preferred", "provider": "NRCan", "grant_type": "non_refundable",
     "agency_details": {"portal_url": "https://nrcan-funding-financement-rncan.canada.ca", "portal_name": "NRCan Funding Portal", "contact_phone": "343-292-6096", "contact_email": "nrcan.eip-pie.rncan@canada.ca", "deadline": "Via targeted calls for proposals - check portal", "processing_time": "120 business days from submission", "how_to_apply": "1. Register on NRCan Funding Portal with GCKey\n2. Monitor current calls for proposals page\n3. Submit Expression of Interest (EOI) during open call\n4. If invited, submit Full Project Proposal (FPP)\n5. Include budget, timeline, partner letters, GHG estimates"}},
    {"name": "Digital Technology Adoption", "description": "Support for small businesses adopting digital technologies including e-commerce, digital marketing, and IT infrastructure. Note: CDAP closed March 2024 - check for successor programs.", "max_amount": 15000, "category": "Digital", "eligibility": "Canadian SME looking to adopt digital tech, 1-499 employees", "provider": "ISED Canada", "grant_type": "non_refundable",
     "agency_details": {"portal_url": "https://ised-isde.canada.ca/site/canada-digital-adoption-program/en", "portal_name": "ISED Innovation Portal", "contact_phone": "1-800-328-6189", "contact_email": "ic.cdap-patn.ic@ised-isde.gc.ca", "deadline": "Original CDAP closed March 2024 - monitor ISED for new programs", "processing_time": "Varies by program", "how_to_apply": "1. Check ISED website for current digital adoption programs\n2. Register on the relevant portal\n3. Work with an approved digital advisor\n4. Submit digital adoption plan\n5. Receive approval and funding"}},
    {"name": "IRAP Innovation Assistance", "description": "NRC Industrial Research Assistance Program for technology innovation and R&D projects. Provides both funding and advisory services.", "max_amount": 100000, "category": "Innovation", "eligibility": "Canadian SME with innovative technology project, incorporated in Canada", "provider": "NRC-IRAP", "grant_type": "refundable",
     "agency_details": {"portal_url": "https://nrc.canada.ca/en/support-technology-innovation", "portal_name": "NRC IRAP Portal", "contact_phone": "1-877-994-4727", "contact_email": "irap.information-information.pari@nrc-cnrc.gc.ca", "deadline": "Rolling intake - apply anytime (fiscal year April 1 - March 31)", "processing_time": "2-6 months from initial contact to approval", "how_to_apply": "1. Contact NRC IRAP at 1-877-994-4727 or register online\n2. Get assigned an Industrial Technology Advisor (ITA)\n3. ITA assesses your project and helps develop proposal\n4. Submit comprehensive proposal through secure IRAP portal\n5. Include technical objectives, milestones, and budget"}},
    {"name": "Health Technology Fund", "description": "Funding for health-tech startups and SMEs developing digital health solutions, medical devices, or health data platforms.", "max_amount": 125000, "category": "Health Tech", "eligibility": "Canadian health-tech company with a live product or prototype", "provider": "Health Canada / PHAC", "grant_type": "non_refundable",
     "agency_details": {"portal_url": "https://fas.hc-phac.canada.ca", "portal_name": "Health Canada Funding Application System (FAS)", "contact_phone": "1-866-225-0709", "contact_email": "fas-sdf@canada.ca", "deadline": "Via targeted calls - check PHAC funding page", "processing_time": "90-180 days", "how_to_apply": "1. Register at Health Canada FAS portal (fas.hc-phac.canada.ca)\n2. Monitor PHAC funding opportunities page for open calls\n3. Prepare application with project description and budget\n4. Submit through the FAS portal by deadline\n5. Also check CIHR grants for health research funding"}},
]

async def seed_all():
    admin_email = os.environ.get("ADMIN_EMAIL", "admin@granthunter.ai")
    admin_password = os.environ.get("ADMIN_PASSWORD", "changeme")
    existing_admin = await db.users.find_one({"email": admin_email})
    if not existing_admin:
        await db.users.insert_one({"email": admin_email, "password_hash": hash_password(admin_password), "name": "Admin", "role": "admin", "business_name": "", "business_description": "", "created_at": datetime.now(timezone.utc).isoformat()})
    elif not verify_password(admin_password, existing_admin["password_hash"]):
        await db.users.update_one({"email": admin_email}, {"$set": {"password_hash": hash_password(admin_password)}})

    if await db.grants.count_documents({}) == 0:
        for g in SEED_GRANTS:
            await db.grants.insert_one({"id": str(uuid.uuid4()), **g, "created_at": datetime.now(timezone.utc).isoformat()})

@app.on_event("startup")
async def startup():
    await db.users.create_index("email", unique=True)
    try:
        init_storage()
        logger.info("Object storage initialized")
    except Exception as e:
        logger.warning(f"Storage init deferred: {e}")
    await seed_all()
    logger.info("GrantGrabber backend started")

app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
